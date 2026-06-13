"""
generate_carousel.py
--------------------
Generates a social media carousel PPTX from a story in aipulse.db.

Usage:
    python generate_carousel.py --story-id 42
    python generate_carousel.py --story-id 42 --output my_post.pptx
    python generate_carousel.py --list            # show available stories
    python generate_carousel.py --story-id 42 --no-images   # skip AI image gen

Requirements:
    pip install python-pptx requests Pillow python-dotenv
"""

import argparse
import json
import os
import re
import sqlite3
import sys
import time
from datetime import datetime
from pathlib import Path

import requests
from dotenv import load_dotenv

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    import lxml.etree as etree
except ImportError:
    print("ERROR: Missing dependencies. Run:\n  pip install python-pptx lxml requests python-dotenv Pillow")
    sys.exit(1)

# ─────────────────────────────────────────────
# Paths
# ─────────────────────────────────────────────

SCRIPT_DIR    = Path(__file__).parent
BACKEND_DIR   = SCRIPT_DIR.parent
DB_PATH       = BACKEND_DIR / "aipulse.db"
TEMPLATE_PATH = SCRIPT_DIR / "carousel_template.pptx"
BRAND_CONFIG  = SCRIPT_DIR / "brand_config.json"
OUTPUT_DIR    = SCRIPT_DIR / "output"

load_dotenv(BACKEND_DIR / ".env")

# ─────────────────────────────────────────────
# API Keys
# ─────────────────────────────────────────────

OPENROUTER_KEY = os.getenv("OPENROUTER_API_KEY")
OPENAI_KEY     = os.getenv("OPENAI_API_KEY")
ANTHROPIC_KEY  = os.getenv("ANTHROPIC_API_KEY")
PIAPI_KEY      = os.getenv("PIAPI_API_KEY")
LLM_MODEL      = os.getenv("LLM_MODEL", "meta-llama/llama-3-70b-instruct:free")

# ─────────────────────────────────────────────
# Text Utilities
# ─────────────────────────────────────────────

def truncate(text, max_chars):
    """Truncate text to max_chars, breaking at word boundary."""
    if not text:
        return ""
    text = text.strip()
    if len(text) <= max_chars:
        return text
    clipped = text[:max_chars].rsplit(' ', 1)[0]
    return clipped.rstrip('.,;:') + "…"


# ─────────────────────────────────────────────
# LLM Helpers
# ─────────────────────────────────────────────

def call_llm(prompt):
    if ANTHROPIC_KEY and not ANTHROPIC_KEY.startswith("your_"):
        return _call_anthropic(prompt)
    api_key = OPENROUTER_KEY or OPENAI_KEY
    if api_key and not api_key.startswith("your_"):
        return _call_openrouter(prompt, api_key)
    print("WARNING: No valid API key found. Using placeholder content.")
    return _placeholder_content()


def _call_anthropic(prompt):
    headers = {
        "x-api-key": ANTHROPIC_KEY,
        "anthropic-version": "2023-06-01",
        "content-type": "application/json",
    }
    payload = {
        "model": "claude-haiku-4-5-20251001",
        "max_tokens": 1500,
        "messages": [{"role": "user", "content": prompt}],
    }
    resp = requests.post("https://api.anthropic.com/v1/messages", headers=headers, json=payload, timeout=30)
    resp.raise_for_status()
    return resp.json()["content"][0]["text"]


def _call_openrouter(prompt, api_key):
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }
    payload = {
        "model": LLM_MODEL,
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.7,
    }
    resp = requests.post("https://openrouter.ai/api/v1/chat/completions", headers=headers, json=payload, timeout=30)
    resp.raise_for_status()
    return resp.json()["choices"][0]["message"]["content"].strip()


def _placeholder_content():
    return json.dumps({
        "post_type": "AI Breakdown",
        "images": {
            "slide3": "Futuristic AI neural network glowing blue circuits on dark background, cinematic lighting, wide landscape",
            "slide4": "Person using laptop with holographic AI interface, professional office setting, wide landscape",
            "avatar": "Minimalist glowing brain icon, flat design, green neon on dark background, square"
        },
        "slide1": {
            "headline_green": "Breaking",
            "headline_white": "News",
            "hook_pill": "What You Need To Know",
            "hook_body_plain": "A major development is reshaping the AI landscape. ",
            "hook_body_highlight": "Here's why it matters",
            "hook_body_plain2": " for your work and business."
        },
        "slide2": {
            "slide_tagline": "The key takeaways",
            "col1": "This release signals a major shift in how AI models handle complex reasoning tasks at scale.",
            "col2": "Open-source alternatives are closing the gap fast, putting pressure on major labs.",
            "col3": "Businesses that adopt early will gain a significant productivity edge over competitors.",
            "headline_green": "What This",
            "headline_white": "Means For You",
            "hook_pill": "Three Things To Know"
        },
        "slide3": {
            "slide_tagline": "The full picture",
            "headline_green": "Here's The",
            "headline_white": "Real Story",
            "left_body_plain": "The announcement caught the industry by surprise. ",
            "left_body_highlight": "Experts call it a landmark moment",
            "left_body_plain2": " in the history of AI development.",
            "right_body_plain": "The ",
            "right_body_highlight": "implications are broad",
            "right_body_plain2": " and felt across every sector."
        },
        "slide4": {
            "slide_tagline": "Take action now",
            "hook_pill": "Your Next Move",
            "headline_green": "How To",
            "headline_white": "Stay Ahead",
            "body_plain": "Start by understanding the core change. ",
            "body_highlight": "Audit your current AI tools and workflows",
            "body_plain2": " to identify where this applies, then act fast."
        },
        "slide5": {
            "closing_word_green": "Stay",
            "closing_word_white": "Sharp",
            "closing_body_plain": "The AI space moves fast — ",
            "closing_body_highlight": "stay curious and keep learning.",
            "closing_body_plain2": " See you next week."
        }
    })


# ─────────────────────────────────────────────
# Content Generation
# ─────────────────────────────────────────────

def generate_slide_content(story, brand):
    category   = brand.get("category_tag", "AI & Tech")
    tagline    = brand.get("brand_tagline", "Stay ahead of AI")
    story_text = story.get("full_text") or story.get("clean_summary") or story.get("summary") or story.get("title", "")
    story_text = story_text[:3000]

    prompt = f"""You are a social media content expert creating a 5-slide carousel post for {category}.

STORY:
Title: {story["title"]}
Source: {story.get("source", "")}
Content: {story_text}

BRAND CONTEXT:
Category: {category}
Tagline: {tagline}

Create carousel slide content. Return ONLY valid JSON with this exact structure:

{{
  "post_type": "3-4 word content label for this specific story (e.g. 'AI Breakthrough', 'Tool Release')",
  "images": {{
    "slide3": "Detailed prompt for a LANDSCAPE photorealistic/cinematic image representing this story. No text overlays. ~20 words.",
    "slide4": "Detailed prompt for a second LANDSCAPE image showing practical impact of this story. Different from slide3. ~20 words.",
    "avatar": "Prompt for a SQUARE minimalist icon representing this story topic. Simple, bold, green neon on dark. ~15 words."
  }},
  "slide1": {{
    "headline_green": "1-2 words GREEN (max 12 chars total)",
    "headline_white": "1-2 words WHITE (max 12 chars total)",
    "hook_pill": "5-7 word pill badge text",
    "hook_body_plain": "Opening sentence plain ~10 words. ",
    "hook_body_highlight": "key phrase bold green 3-5 words",
    "hook_body_plain2": " closing plain ~8 words."
  }},
  "slide2": {{
    "slide_tagline": "5-6 word right-nav tagline",
    "col1": "First insight ~15 words punchy.",
    "col2": "Second insight ~15 words punchy.",
    "col3": "Third insight ~15 words punchy.",
    "headline_green": "2 word green",
    "headline_white": "3-4 word white",
    "hook_pill": "4-5 word pill"
  }},
  "slide3": {{
    "slide_tagline": "5-6 word right-nav tagline",
    "headline_green": "1-2 word green",
    "headline_white": "3-5 word white",
    "left_body_plain": "Left col opening ~10 words. ",
    "left_body_highlight": "key phrase 3-4 words",
    "left_body_plain2": " closing ~8 words.",
    "right_body_plain": "Right col opening ~5 words ",
    "right_body_highlight": "key phrase 2-3 words",
    "right_body_plain2": " closing ~6 words."
  }},
  "slide4": {{
    "slide_tagline": "5-6 word right-nav tagline",
    "hook_pill": "4-5 word pill",
    "headline_green": "1-2 word green",
    "headline_white": "3-4 word white",
    "body_plain": "Opening body ~10 words. ",
    "body_highlight": "key actionable phrase 4-6 words",
    "body_plain2": " closing ~10 words."
  }},
  "slide5": {{
    "closing_word_green": "FIRST word only (e.g. Stay)",
    "closing_word_white": "SECOND word only (e.g. Sharp)",
    "closing_body_plain": "Opening ~8 words — ",
    "closing_body_highlight": "key inspiring phrase 4-6 words.",
    "closing_body_plain2": " Short sign-off ~5 words."
  }}
}}

Return only the JSON. No explanation, no markdown fences."""

    print("  Calling LLM to generate slide content...")
    raw = call_llm(prompt)
    raw = raw.strip()
    if raw.startswith("```"):
        raw = re.sub(r"^```[a-z]*\n?", "", raw)
        raw = re.sub(r"\n?```$", "", raw)
    try:
        return json.loads(raw)
    except json.JSONDecodeError as e:
        print(f"  WARNING: Could not parse LLM response ({e}). Using placeholder.")
        return json.loads(_placeholder_content())


# ─────────────────────────────────────────────
# PiAPI Image Generation
# ─────────────────────────────────────────────

def _submit_piapi_task(prompt, width, height):
    headers = {
        "X-API-Key": PIAPI_KEY,
        "Content-Type": "application/json",
    }
    payload = {
        "model": "Qubico/flux1-schnell",
        "task_type": "txt2img",
        "input": {"prompt": prompt, "width": width, "height": height},
    }
    resp = requests.post("https://api.piapi.ai/api/v1/task", headers=headers, json=payload, timeout=30)
    resp.raise_for_status()
    data = resp.json()
    if data.get("code") != 200:
        raise RuntimeError(f"PiAPI submit error: {data}")
    return data["data"]["task_id"]


def _poll_piapi_task(task_id, timeout=120):
    headers = {"X-API-Key": PIAPI_KEY}
    deadline = time.time() + timeout
    while time.time() < deadline:
        resp = requests.get(
            f"https://api.piapi.ai/api/v1/task/{task_id}",
            headers=headers,
            timeout=15,
        )
        resp.raise_for_status()
        data = resp.json()
        if data.get("code") != 200:
            raise RuntimeError(f"PiAPI poll error: {data}")
        task_data = data["data"]
        status = task_data.get("status", "")
        if status == "completed":
            output = task_data.get("output", {})
            if "image_url" in output:
                return output["image_url"]
            if "images" in output and output["images"]:
                return output["images"][0]
            raise RuntimeError(f"Task completed but no image URL: {output}")
        if status in ("failed", "error"):
            raise RuntimeError(f"Task {task_id} failed: {task_data.get('error', status)}")
        time.sleep(3)
    raise TimeoutError(f"Task {task_id} timed out after {timeout}s")


def generate_images_piapi(content):
    if not PIAPI_KEY or PIAPI_KEY.startswith("your_"):
        print("  Skipping AI image generation (no PIAPI_API_KEY set).")
        return {}

    images_spec = content.get("images", {})
    tasks_to_submit = [
        ("slide3", images_spec.get("slide3", ""), 1024, 512),
        ("slide4", images_spec.get("slide4", ""), 1024, 512),
        ("avatar", images_spec.get("avatar", ""), 512, 512),
    ]

    task_ids = {}
    for key, prompt, w, h in tasks_to_submit:
        if not prompt:
            continue
        try:
            print(f"  Submitting [{key}]: {prompt[:70]}...")
            task_ids[key] = _submit_piapi_task(prompt, w, h)
            print(f"    task_id: {task_ids[key]}")
        except Exception as e:
            print(f"  WARNING: Submit [{key}] failed: {e}")

    image_urls = {}
    for key, task_id in task_ids.items():
        try:
            print(f"  Waiting for [{key}]...")
            url = _poll_piapi_task(task_id)
            image_urls[key] = url
            print(f"    OK: {url[:90]}")
        except Exception as e:
            print(f"  WARNING: [{key}] image failed: {e}")

    return image_urls


# ─────────────────────────────────────────────
# PPTX Helpers
# ─────────────────────────────────────────────

def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
        if shape.shape_type == 6:
            for child in shape.shapes:
                if child.name == name:
                    return child
    return None


def replace_simple(slide, shape_name, new_text):
    shape = find_shape(slide, shape_name)
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    para = tf.paragraphs[0]
    if not para.runs:
        return
    para.runs[0].text = new_text
    for run in para.runs[1:]:
        run.text = ""
    for p in tf.paragraphs[1:]:
        for run in p.runs:
            run.text = ""


def replace_headline(slide, shape_name, green_part, white_part):
    shape = find_shape(slide, shape_name)
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    para = tf.paragraphs[0]
    runs = para.runs
    if len(runs) >= 2:
        runs[0].text = green_part
        runs[1].text = (" " + white_part) if not white_part.startswith(" ") else white_part
        for run in runs[2:]:
            run.text = ""
    elif len(runs) == 1:
        runs[0].text = green_part + " " + white_part
    for p in tf.paragraphs[1:]:
        for run in p.runs:
            run.text = ""


def replace_body_with_highlight(slide, shape_name, plain1, highlight, plain2):
    shape = find_shape(slide, shape_name)
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    para = tf.paragraphs[0]
    runs = para.runs
    if len(runs) >= 3:
        runs[0].text = plain1
        runs[1].text = highlight
        runs[2].text = plain2
        for run in runs[3:]:
            run.text = ""
    elif len(runs) == 2:
        runs[0].text = plain1
        runs[1].text = highlight + plain2
    elif len(runs) == 1:
        runs[0].text = plain1 + highlight + plain2
    for p in tf.paragraphs[1:]:
        for run in p.runs:
            run.text = ""


def replace_pill_text(slide, group_name, new_text):
    for shape in slide.shapes:
        if shape.name == group_name and shape.shape_type == 6:
            for child in shape.shapes:
                if child.has_text_frame and child.text_frame.text.strip():
                    tf = child.text_frame
                    if tf.paragraphs and tf.paragraphs[0].runs:
                        tf.paragraphs[0].runs[0].text = new_text
                        for run in tf.paragraphs[0].runs[1:]:
                            run.text = ""
                    return


def _replace_blipfill(slide, group_name, image_data):
    for shape in slide.shapes:
        if shape.name == group_name and shape.shape_type == 6:
            for child in shape.shapes:
                xml_str = etree.tostring(child._element).decode()
                if "blipFill" not in xml_str:
                    continue
                rids = re.findall(r'r:embed="([^"]+)"', xml_str)
                if not rids:
                    continue
                try:
                    slide.part.rels[rids[0]].target_part._blob = image_data
                    return True
                except Exception as e:
                    print(f"  WARNING: blipFill replace failed [{group_name}]: {e}")
                    return False
    return False


def replace_image_in_shape(slide, group_name, image_url):
    try:
        resp = requests.get(image_url, timeout=20)
        resp.raise_for_status()
        if _replace_blipfill(slide, group_name, resp.content):
            print(f"  Content image replaced [{group_name}]")
        else:
            print(f"  WARNING: blipFill target not found [{group_name}]")
    except Exception as e:
        print(f"  WARNING: Image download failed [{group_name}]: {e}")


def replace_avatar_image(slide, image_data):
    """Replace the avatar circle (bottom-left) by position: top 13-15\", left < 2\"."""
    for shape in slide.shapes:
        if shape.shape_type != 6:
            continue
        top_in  = shape.top  / 914400
        left_in = shape.left / 914400
        if not (13.0 < top_in < 15.0 and left_in < 2.0):
            continue
        for child in shape.shapes:
            xml_str = etree.tostring(child._element).decode()
            if "blipFill" not in xml_str:
                continue
            rids = re.findall(r'r:embed="([^"]+)"', xml_str)
            if not rids:
                continue
            try:
                slide.part.rels[rids[0]].target_part._blob = image_data
                return True
            except Exception as e:
                print(f"  WARNING: Avatar replace failed: {e}")
                return False
    return False


# ─────────────────────────────────────────────
# Database Helpers
# ─────────────────────────────────────────────

def get_story(story_id):
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cur = conn.cursor()
    cur.execute("SELECT * FROM stories WHERE id = ?", (story_id,))
    row = cur.fetchone()
    conn.close()
    if not row:
        print(f"ERROR: No story found with id={story_id}")
        sys.exit(1)
    return dict(row)


def list_stories(limit=20):
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cur = conn.cursor()
    cur.execute("""
        SELECT id, title, source, value_score, status, created_at
        FROM stories ORDER BY id DESC LIMIT ?
    """, (limit,))
    rows = cur.fetchall()
    conn.close()
    return [dict(r) for r in rows]


def get_story_images(story):
    raw = story.get("images_json") or "[]"
    try:
        data = json.loads(raw)
        if isinstance(data, list):
            return [img for img in data if isinstance(img, str) and img.startswith("http")]
        return []
    except Exception:
        return []


# ─────────────────────────────────────────────
# Carousel Builder
# ─────────────────────────────────────────────

def build_carousel(story, brand, content, output_path, image_urls):
    """Populate the template PPTX with content and images, then save."""
    prs    = Presentation(str(TEMPLATE_PATH))
    slides = prs.slides

    # Brand values
    author        = truncate(brand.get("author_name", "Your Name"), 30)
    cat_tag       = truncate(brand.get("category_tag", "AI & Tech"), 20)
    brand_tagline = truncate(brand.get("brand_tagline", "Stay ahead of AI"), 35)
    pill_tag      = truncate(brand.get("pill_tagline", "Stay Ahead. Stay Sharp"), 40)
    closing_cta   = truncate(brand.get("closing_cta", "Follow for daily AI insights"), 40)
    contact_pill1 = truncate(brand.get("contact_pill_1", pill_tag), 40)
    contact_pill2 = truncate(brand.get("contact_pill_2", pill_tag), 40)

    # Per-story center nav label from LLM, fallback to brand config
    post_type = truncate(
        content.get("post_type") or brand.get("post_type_label", "Weekly Carousel"),
        25
    )

    s1 = content.get("slide1", {})
    s2 = content.get("slide2", {})
    s3 = content.get("slide3", {})
    s4 = content.get("slide4", {})
    s5 = content.get("slide5", {})

    # SLIDE 1: Cover
    sl1 = slides[0]
    replace_simple(sl1,    "TextBox 19", cat_tag)
    replace_simple(sl1,    "TextBox 20", post_type)
    replace_simple(sl1,    "TextBox 21", brand_tagline)
    replace_simple(sl1,    "TextBox 18", author)
    replace_headline(sl1,  "TextBox 16",
                     truncate(s1.get("headline_green", "Breaking"), 14),
                     truncate(s1.get("headline_white", "News"), 14))
    replace_pill_text(sl1, "Group 10", truncate(s1.get("hook_pill", pill_tag), 40))
    replace_simple(sl1,    "TextBox 22", truncate(s1.get("hook_pill", pill_tag), 40))
    replace_body_with_highlight(sl1, "TextBox 17",
        truncate(s1.get("hook_body_plain", ""), 120),
        truncate(s1.get("hook_body_highlight", ""), 60),
        truncate(s1.get("hook_body_plain2", ""), 120))

    # SLIDE 2: 3-Column
    sl2 = slides[1]
    replace_simple(sl2,    "TextBox 9",  cat_tag)
    replace_simple(sl2,    "TextBox 10", post_type)
    replace_simple(sl2,    "TextBox 11", truncate(s2.get("slide_tagline", brand_tagline), 40))
    replace_simple(sl2,    "TextBox 8",  author)
    replace_headline(sl2,  "TextBox 12",
                     truncate(s2.get("headline_green", "Key"), 14),
                     truncate(s2.get("headline_white", "Insights"), 16))
    replace_pill_text(sl2, "Group 10", truncate(s2.get("hook_pill", pill_tag), 40))
    replace_body_with_highlight(sl2, "TextBox 13", truncate(s2.get("col1", ""), 100), "", "")
    replace_body_with_highlight(sl2, "TextBox 14", truncate(s2.get("col2", ""), 100), "", "")
    replace_body_with_highlight(sl2, "TextBox 15", truncate(s2.get("col3", ""), 100), "", "")

    # SLIDE 3: Image + 2-Column
    sl3 = slides[2]
    replace_simple(sl3,    "TextBox 13", cat_tag)
    replace_simple(sl3,    "TextBox 14", post_type)
    replace_simple(sl3,    "TextBox 15", truncate(s3.get("slide_tagline", brand_tagline), 40))
    replace_simple(sl3,    "TextBox 12", author)
    replace_headline(sl3,  "TextBox 16",
                     truncate(s3.get("headline_green", "The"), 14),
                     truncate(s3.get("headline_white", "Full Picture"), 16))
    replace_body_with_highlight(sl3, "TextBox 17",
        truncate(s3.get("left_body_plain", ""), 120),
        truncate(s3.get("left_body_highlight", ""), 60),
        truncate(s3.get("left_body_plain2", ""), 120))
    replace_body_with_highlight(sl3, "TextBox 18",
        truncate(s3.get("right_body_plain", ""), 100),
        truncate(s3.get("right_body_highlight", ""), 50),
        truncate(s3.get("right_body_plain2", ""), 100))
    if "slide3" in image_urls:
        replace_image_in_shape(sl3, "Group 8", image_urls["slide3"])

    # SLIDE 4: Image + Headline
    sl4 = slides[3]
    replace_simple(sl4,    "TextBox 19", cat_tag)
    replace_simple(sl4,    "TextBox 20", post_type)
    replace_simple(sl4,    "TextBox 21", truncate(s4.get("slide_tagline", brand_tagline), 40))
    replace_simple(sl4,    "TextBox 18", author)
    replace_headline(sl4,  "TextBox 22",
                     truncate(s4.get("headline_green", "Take"), 14),
                     truncate(s4.get("headline_white", "Action Now"), 16))
    replace_pill_text(sl4, "Group 10", truncate(s4.get("hook_pill", pill_tag), 40))
    replace_body_with_highlight(sl4, "TextBox 23",
        truncate(s4.get("body_plain", ""), 120),
        truncate(s4.get("body_highlight", ""), 60),
        truncate(s4.get("body_plain2", ""), 120))
    if "slide4" in image_urls:
        replace_image_in_shape(sl4, "Group 14", image_urls["slide4"])

    # SLIDE 5: Closing
    sl5 = slides[4]
    replace_simple(sl5,    "TextBox 9",  cat_tag)
    replace_simple(sl5,    "TextBox 10", post_type)
    replace_simple(sl5,    "TextBox 11", closing_cta)   # brand CTA, not "Thank you for reading"
    replace_simple(sl5,    "TextBox 8",  author)
    replace_headline(sl5,  "TextBox 12",
                     truncate(s5.get("closing_word_green", "Stay"), 10),
                     truncate(s5.get("closing_word_white", "Sharp"), 12))
    replace_body_with_highlight(sl5, "TextBox 16",
        truncate(s5.get("closing_body_plain", ""), 120),
        truncate(s5.get("closing_body_highlight", ""), 60),
        truncate(s5.get("closing_body_plain2", ""), 120))

    # Slide 5 two pill badges: left=contact_pill1, right=contact_pill2
    pill_shapes = sorted(
        [s for s in sl5.shapes if s.name == "Group 10" and s.shape_type == 6],
        key=lambda s: s.left,
    )
    for i, pill_shape in enumerate(pill_shapes):
        text = contact_pill1 if i == 0 else contact_pill2
        for child in pill_shape.shapes:
            if child.has_text_frame and child.text_frame.text.strip():
                tf = child.text_frame
                if tf.paragraphs and tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = text
                    for run in tf.paragraphs[0].runs[1:]:
                        run.text = ""
                break

    # Avatar image applied to all slides
    if "avatar" in image_urls and image_urls["avatar"]:
        try:
            print("  Downloading avatar image...")
            resp = requests.get(image_urls["avatar"], timeout=20)
            resp.raise_for_status()
            avatar_data = resp.content
            count = sum(1 for slide in slides if replace_avatar_image(slide, avatar_data))
            print(f"  Avatar applied to {count}/{len(prs.slides)} slides")
        except Exception as e:
            print(f"  WARNING: Avatar image failed: {e}")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"\nCarousel saved: {output_path}")


# ─────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Generate a carousel PPTX from a story in aipulse.db"
    )
    parser.add_argument("--story-id",  type=int,  help="Story ID from aipulse.db")
    parser.add_argument("--output",    type=str,  help="Output PPTX path (optional)")
    parser.add_argument("--list",      action="store_true", help="List recent stories")
    parser.add_argument("--no-images", action="store_true", help="Skip AI image generation")
    args = parser.parse_args()

    if args.list:
        print(f"\n{'ID':<6} {'Score':<7} {'Source':<20} {'Status':<12} Title")
        print("-" * 90)
        for s in list_stories():
            title = (s["title"] or "")[:55]
            print(f"{s['id']:<6} {(s['value_score'] or 0):<7} {(s['source'] or ''):<20} {(s['status'] or ''):<12} {title}")
        return

    if not args.story_id:
        parser.print_help()
        return

    if not BRAND_CONFIG.exists():
        print(f"ERROR: brand_config.json not found at {BRAND_CONFIG}")
        sys.exit(1)
    with open(BRAND_CONFIG) as f:
        brand = json.load(f)

    print(f"\nLoading story #{args.story_id}...")
    story = get_story(args.story_id)
    print(f"  Title:  {story['title']}")
    print(f"  Source: {story.get('source', 'unknown')}")

    print("\nGenerating slide content...")
    content = generate_slide_content(story, brand)

    image_urls = {}
    if not args.no_images:
        print("\nGenerating AI images via PiAPI FLUX.1...")
        image_urls = generate_images_piapi(content)
        if not image_urls:
            print("  No images generated — continuing without.")
    else:
        print("\nSkipping AI image generation (--no-images).")

    if args.output:
        output_path = Path(args.output)
    else:
        safe_title = re.sub(r'[^\w\s-]', '', story['title'][:40]).strip().replace(' ', '_')
        date_str   = datetime.now().strftime("%Y%m%d")
        output_path = OUTPUT_DIR / f"carousel_{date_str}_{safe_title}.pptx"

    print("\nBuilding PPTX...")
    build_carousel(story, brand, content, output_path, image_urls)

    print(f"\n{'='*55}")
    print(f"  Story:   {story['title'][:55]}")
    print(f"  Images:  {len(image_urls)} generated")
    print(f"  Output:  {output_path}")
    print(f"{'='*55}\n")


if __name__ == "__main__":
    main()
