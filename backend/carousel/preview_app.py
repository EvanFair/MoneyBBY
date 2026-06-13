"""
preview_app.py  —  Carousel Preview App
Run with:  python preview_app.py
Opens http://localhost:5050 automatically.
"""

import argparse, json, os, re, sqlite3, subprocess, sys, threading, webbrowser
from datetime import datetime
from pathlib import Path
from flask import Flask, jsonify, request

SCRIPT_DIR  = Path(__file__).parent
BACKEND_DIR = SCRIPT_DIR.parent
DB_PATH     = BACKEND_DIR / "aipulse.db"
OUTPUT_DIR  = SCRIPT_DIR / "output"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

app = Flask(__name__, static_folder=None)

# ── DB ─────────────────────────────────────────────────────────────────────────

def get_curated_stories():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cur = conn.cursor()
    cur.execute("""
        SELECT id, title, source, value_score, status, created_at,
               COALESCE(summary, clean_summary, '') as summary
        FROM stories
        WHERE status = 'approved' OR COALESCE(value_score, 0) > 5
        ORDER BY COALESCE(value_score, 0) DESC, created_at DESC
        LIMIT 100
    """)
    rows = cur.fetchall()
    conn.close()
    return [dict(r) for r in rows]

def get_story(story_id):
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cur = conn.cursor()
    cur.execute("SELECT * FROM stories WHERE id=?", (story_id,))
    row = cur.fetchone()
    conn.close()
    return dict(row) if row else None

# ── Slide data extraction ───────────────────────────────────────────────────────

def extract_slides(pptx_path):
    from pptx import Presentation
    prs = Presentation(str(pptx_path))
    results = []

    for idx, slide in enumerate(prs.slides):

        def t(name):
            for sh in slide.shapes:
                if sh.name == name and sh.has_text_frame:
                    return sh.text_frame.text.strip()
                if sh.shape_type == 6:
                    for c in sh.shapes:
                        if c.name == name and c.has_text_frame:
                            return c.text_frame.text.strip()
            return ""

        def hl(name):
            for sh in slide.shapes:
                if sh.name == name and sh.has_text_frame and sh.text_frame.paragraphs:
                    runs = sh.text_frame.paragraphs[0].runs
                    if len(runs) >= 2:
                        return runs[0].text.strip(), runs[1].text.strip()
                    if len(runs) == 1:
                        parts = runs[0].text.strip().split(" ", 1)
                        return parts[0], parts[1] if len(parts) > 1 else ""
            return "", ""

        def bd(name):
            for sh in slide.shapes:
                if sh.name == name and sh.has_text_frame and sh.text_frame.paragraphs:
                    runs = sh.text_frame.paragraphs[0].runs
                    if len(runs) >= 3: return runs[0].text, runs[1].text, runs[2].text
                    if len(runs) == 2: return runs[0].text, runs[1].text, ""
                    if len(runs) == 1: return runs[0].text, "", ""
            return "", "", ""

        def pill(gname):
            for sh in slide.shapes:
                if sh.name == gname and sh.shape_type == 6:
                    for c in sh.shapes:
                        if c.has_text_frame and c.text_frame.text.strip():
                            return c.text_frame.text.strip()
            return ""

        if idx == 0:
            hg, hw = hl("TextBox 16")
            b1, bh, b2 = bd("TextBox 17")
            results.append({"type":"cover",
                "nav": [t("TextBox 19"), t("TextBox 20"), t("TextBox 21")],
                "author": t("TextBox 18"),
                "headline_green": hg, "headline_white": hw,
                "body": [b1, bh, b2], "pill": pill("Group 10")})

        elif idx == 1:
            hg, hw = hl("TextBox 12")
            results.append({"type":"3col",
                "nav": [t("TextBox 9"), t("TextBox 10"), t("TextBox 11")],
                "author": t("TextBox 8"),
                "headline_green": hg, "headline_white": hw,
                "cols": [t("TextBox 13"), t("TextBox 14"), t("TextBox 15")],
                "pill": pill("Group 10")})

        elif idx == 2:
            hg, hw = hl("TextBox 16")
            b1l, bhl, b2l = bd("TextBox 17")
            b1r, bhr, b2r = bd("TextBox 18")
            results.append({"type":"image_2col",
                "nav": [t("TextBox 13"), t("TextBox 14"), t("TextBox 15")],
                "author": t("TextBox 12"),
                "headline_green": hg, "headline_white": hw,
                "body_left": [b1l, bhl, b2l], "body_right": [b1r, bhr, b2r]})

        elif idx == 3:
            hg, hw = hl("TextBox 22")
            b1, bh, b2 = bd("TextBox 23")
            results.append({"type":"cta",
                "nav": [t("TextBox 19"), t("TextBox 20"), t("TextBox 21")],
                "author": t("TextBox 18"),
                "headline_green": hg, "headline_white": hw,
                "body": [b1, bh, b2], "pill": pill("Group 10")})

        elif idx == 4:
            hg, hw = hl("TextBox 12")
            b1, bh, b2 = bd("TextBox 16")
            sorted_pills = sorted(
                [s for s in slide.shapes if s.name == "Group 10" and s.shape_type == 6],
                key=lambda s: s.left)
            ptexts = []
            for ps in sorted_pills:
                for c in ps.shapes:
                    if c.has_text_frame and c.text_frame.text.strip():
                        ptexts.append(c.text_frame.text.strip()); break
            results.append({"type":"closing",
                "nav": [t("TextBox 9"), t("TextBox 10"), t("TextBox 11")],
                "author": t("TextBox 8"),
                "closing_green": hg, "closing_white": hw,
                "body": [b1, bh, b2], "pills": ptexts})

    return results

# ── Routes ─────────────────────────────────────────────────────────────────────

@app.route("/")
def index():
    return HTML

@app.route("/api/stories")
def api_stories():
    try:
        return jsonify({"ok": True, "stories": get_curated_stories()})
    except Exception as e:
        return jsonify({"ok": False, "error": str(e)}), 500

@app.route("/api/generate/<int:story_id>", methods=["POST"])
def api_generate(story_id):
    story = get_story(story_id)
    if not story:
        return jsonify({"ok": False, "error": f"Story {story_id} not found"}), 404

    safe = re.sub(r"[^\w\s-]", "", story["title"][:40]).strip().replace(" ", "_")
    out  = OUTPUT_DIR / f"carousel_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{safe}.pptx"

    try:
        r = subprocess.run(
            [sys.executable, str(SCRIPT_DIR / "generate_carousel.py"),
             "--story-id", str(story_id), "--output", str(out), "--no-images"],
            capture_output=True, text=True, timeout=120, cwd=str(SCRIPT_DIR))
        if r.returncode != 0:
            return jsonify({"ok": False, "error": r.stderr or r.stdout or "Generation failed"}), 500
    except subprocess.TimeoutExpired:
        return jsonify({"ok": False, "error": "Timed out after 2 minutes"}), 500
    except Exception as e:
        return jsonify({"ok": False, "error": str(e)}), 500

    if not out.exists():
        return jsonify({"ok": False, "error": "PPTX not created — add an API key to .env"}), 500

    return jsonify({"ok": True, "path": str(out), "filename": out.name,
                    "slides": extract_slides(out), "log": r.stdout})

@app.route("/api/open", methods=["POST"])
def api_open():
    path = (request.json or {}).get("path", "")
    if not path or not Path(path).exists():
        return jsonify({"ok": False, "error": "File not found"}), 404
    try:
        os.startfile(path)
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"ok": False, "error": str(e)}), 500

# ── HTML ───────────────────────────────────────────────────────────────────────

HTML = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<title>Carousel Generator</title>
<style>
*{box-sizing:border-box;margin:0;padding:0}
body{background:#0a0a0a;color:#ddd;font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',sans-serif;display:flex;flex-direction:column;height:100vh;overflow:hidden}

/* Header */
.hdr{background:#111;border-bottom:1px solid #1e1e1e;padding:12px 20px;display:flex;align-items:center;gap:12px;flex-shrink:0}
.hdr-logo{background:#C4FF00;color:#000;font-weight:900;font-size:13px;width:26px;height:26px;border-radius:6px;display:flex;align-items:center;justify-content:center}
.hdr h1{font-size:15px;font-weight:600;color:#fff}
.hdr-sub{margin-left:auto;font-size:11px;color:#444}

/* Layout */
.body{display:flex;flex:1;overflow:hidden}

/* Left */
.left{width:340px;flex-shrink:0;border-right:1px solid #1a1a1a;display:flex;flex-direction:column;overflow:hidden}
.left-top{padding:12px 14px 8px;border-bottom:1px solid #1a1a1a;display:flex;align-items:center;justify-content:space-between}
.left-top span{font-size:10px;font-weight:700;color:#444;text-transform:uppercase;letter-spacing:.8px}
.cnt-badge{background:#161616;color:#C4FF00;font-size:11px;font-weight:700;padding:2px 8px;border-radius:10px}
.stories{flex:1;overflow-y:auto;padding:6px}
.stories::-webkit-scrollbar{width:3px}
.stories::-webkit-scrollbar-thumb{background:#222;border-radius:2px}

.sc{padding:10px 12px;border-radius:8px;cursor:pointer;margin-bottom:3px;border:1px solid transparent;transition:all .12s}
.sc:hover{background:#141414;border-color:#222}
.sc.active{background:#111a00;border-color:#C4FF00}
.sc-title{font-size:12px;font-weight:500;color:#bbb;line-height:1.4;margin-bottom:5px;display:-webkit-box;-webkit-line-clamp:2;-webkit-box-orient:vertical;overflow:hidden}
.sc.active .sc-title{color:#fff}
.sc-meta{display:flex;align-items:center;gap:6px}
.sc-src{font-size:10px;color:#3a3a3a;background:#161616;padding:2px 6px;border-radius:4px}
.sc-score{font-size:11px;font-weight:700;color:#C4FF00}
.sc-status{font-size:10px;padding:1px 5px;border-radius:3px;font-weight:600}
.approved{background:#0a1e00;color:#6a9a20}
.other{background:#161616;color:#444}

/* Right */
.right{flex:1;display:flex;flex-direction:column;overflow:hidden}
.right-hdr{padding:12px 20px;border-bottom:1px solid #1a1a1a;display:flex;align-items:center;gap:10px;min-height:50px;flex-shrink:0}
.right-title{font-size:13px;font-weight:600;color:#fff;flex:1;white-space:nowrap;overflow:hidden;text-overflow:ellipsis}
.slides-wrap{flex:1;overflow-y:auto;padding:20px 24px;display:flex;flex-direction:column;gap:20px}
.slides-wrap::-webkit-scrollbar{width:4px}
.slides-wrap::-webkit-scrollbar-thumb{background:#222;border-radius:2px}

/* Empty / loading */
.empty{display:flex;flex-direction:column;align-items:center;justify-content:center;height:100%;color:#2a2a2a;gap:10px}
.empty .arrow{font-size:32px}
.loading{display:flex;align-items:center;justify-content:center;flex-direction:column;height:200px;gap:14px}
.spin{width:28px;height:28px;border:3px solid #1e1e1e;border-top-color:#C4FF00;border-radius:50%;animation:spin .7s linear infinite}
@keyframes spin{to{transform:rotate(360deg)}}
.loading-msg{font-size:13px;color:#444}

/* Buttons */
.btn{display:inline-flex;align-items:center;gap:6px;padding:6px 14px;border-radius:6px;font-size:12px;font-weight:600;cursor:pointer;border:none;transition:all .12s}
.btn-open{background:#1a1a1a;color:#aaa;border:1px solid #252525}
.btn-open:hover{background:#222;color:#fff}

/* ── Slide mock-up cards ── */
.slide-card{background:#111;border:1px solid #1e1e1e;border-radius:12px;overflow:hidden}
.slide-label{padding:8px 14px;background:#0d0d0d;border-bottom:1px solid #1a1a1a;display:flex;align-items:center;justify-content:space-between}
.slide-label span{font-size:10px;font-weight:700;color:#C4FF00;text-transform:uppercase;letter-spacing:.8px}
.slide-label .slide-n{font-size:10px;color:#333;font-weight:400}

/* Inner mock */
.mock{background:#0d0d0d;padding:0;display:flex;flex-direction:column;min-height:220px}

.sm-nav{display:flex;justify-content:space-between;align-items:center;padding:8px 14px;border-bottom:1px solid #141414}
.sm-nav span{font-size:9px;color:#333}
.sm-nav span:nth-child(2){color:#555;font-weight:600}

.sm-body{padding:14px;flex:1;display:flex;flex-direction:column;gap:8px}
.sm-foot{padding:7px 14px;border-top:1px solid #141414}
.sm-foot span{font-size:9px;color:#333}

/* Slide content elements */
.s-pill{display:inline-block;background:#1a2a00;border:1px solid #2a4a00;color:#C4FF00;font-size:9px;font-weight:700;padding:3px 8px;border-radius:20px;margin-bottom:4px;max-width:100%;white-space:nowrap;overflow:hidden;text-overflow:ellipsis}
.s-hl{font-size:18px;font-weight:800;line-height:1.2;letter-spacing:-0.5px;margin-bottom:6px}
.s-hl .g{color:#C4FF00}
.s-hl .w{color:#fff}
.s-body{font-size:11px;color:#888;line-height:1.5}
.s-body mark{background:transparent;color:#C4FF00;font-weight:600}
.s-cols{display:grid;grid-template-columns:1fr 1fr 1fr;gap:6px;margin-top:4px}
.s-col{background:#111;border:1px solid #1a1a1a;border-radius:5px;padding:6px 7px;font-size:10px;color:#777;line-height:1.4}
.s-two{display:grid;grid-template-columns:1fr 1fr;gap:8px}
.s-img{height:50px;background:#111;border:1px dashed #1e1e1e;border-radius:5px;display:flex;align-items:center;justify-content:center;font-size:9px;color:#2a2a2a;margin-bottom:6px}
.s-closing{font-size:32px;font-weight:900;line-height:1;letter-spacing:-1px;margin-bottom:6px}
.s-closing .g{color:#C4FF00}
.s-closing .w{color:#fff}
.s-pills-row{display:flex;gap:6px;margin-top:6px;flex-wrap:wrap}
</style>
</head>
<body>

<div class="hdr">
  <div class="hdr-logo">C</div>
  <h1>Carousel Generator</h1>
  <span class="hdr-sub">Click a story to generate</span>
</div>

<div class="body">
  <div class="left">
    <div class="left-top">
      <span>Curated Stories</span>
      <span class="cnt-badge" id="cnt">—</span>
    </div>
    <div class="stories" id="list"><div class="loading"><div class="spin"></div></div></div>
  </div>

  <div class="right">
    <div class="right-hdr" id="rhdr">
      <span class="right-title" id="rtitle" style="color:#2a2a2a">← Select a story to generate a carousel</span>
    </div>
    <div class="slides-wrap" id="slides">
      <div class="empty"><div class="arrow">↙</div><p>Pick any story on the left</p></div>
    </div>
  </div>
</div>

<script>
let busy = false;

function esc(s){return String(s||'').replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;')}

async function loadStories(){
  const r = await fetch('/api/stories');
  const d = await r.json();
  if(!d.ok){document.getElementById('list').innerHTML=`<p style="padding:16px;color:#f55;font-size:12px">Error: ${esc(d.error)}</p>`;return}
  document.getElementById('cnt').textContent = d.stories.length;
  if(!d.stories.length){
    document.getElementById('list').innerHTML='<p style="padding:20px;color:#333;font-size:12px;text-align:center">No curated stories yet.<br>Need status=approved or value_score>5.</p>';
    return;
  }
  document.getElementById('list').innerHTML = d.stories.map(s=>{
    const score = s.value_score ? parseFloat(s.value_score).toFixed(1) : '—';
    const src   = (s.source||'').replace(/_/g,' ');
    const scls  = s.status==='approved'?'approved':'other';
    return `<div class="sc" data-id="${s.id}" data-title="${esc(s.title)}" onclick="pick(parseInt(this.dataset.id),this,this.dataset.title)">
      <div class="sc-title">${esc(s.title)}</div>
      <div class="sc-meta">
        <span class="sc-src">${esc(src)}</span>
        <span class="sc-score">★ ${score}</span>
        <span class="sc-status ${scls}">${esc(s.status||'scored')}</span>
      </div>
    </div>`;
  }).join('');
}

async function pick(id, el, title){
  if(busy) return;
  document.querySelectorAll('.sc').forEach(c=>c.classList.remove('active'));
  el.classList.add('active');
  document.getElementById('rtitle').style.color='';
  document.getElementById('rtitle').textContent = title;
  document.getElementById('rhdr').innerHTML = `<span class="right-title">${title}</span>`;
  showLoading();
  busy = true;

  try{
    const r   = await fetch('/api/generate/'+id, {method:'POST'});
    const d   = await r.json();
    if(!d.ok) throw new Error(d.error);
    showSlides(d);
  }catch(e){
    document.getElementById('slides').innerHTML =
      `<div style="padding:24px;color:#f76;font-size:13px;background:#1a0800;border-radius:8px;border:1px solid #4a1800">
        ⚠ ${esc(e.message)}
      </div>`;
  }
  busy = false;
}

function showLoading(){
  document.getElementById('slides').innerHTML =
    `<div class="loading"><div class="spin"></div><div class="loading-msg">Generating slides…</div></div>`;
}

function showSlides(d){
  const names=['Cover','3-Column','Image + Copy','Call To Action','Closing'];
  let html = `<div style="display:flex;align-items:center;gap:10px;margin-bottom:4px">
    <span style="font-size:12px;color:#C4FF00;font-weight:600">✓ ${esc(d.filename)}</span>
    <button class="btn btn-open" data-path="${esc(d.path)}" onclick="openFile(this.dataset.path)">Open in PowerPoint ↗</button>
  </div>`;

  html += d.slides.map((s,i)=>`
    <div class="slide-card">
      <div class="slide-label">
        <span>Slide ${i+1} · ${names[i]||''}</span>
        <span class="slide-n">${esc(s.type)}</span>
      </div>
      <div class="mock">
        <div class="sm-nav">
          <span>${esc(s.nav&&s.nav[0])}</span>
          <span>${esc(s.nav&&s.nav[1])}</span>
          <span>${esc(s.nav&&s.nav[2])}</span>
        </div>
        <div class="sm-body">${renderSlideBody(s)}</div>
        <div class="sm-foot"><span>${esc(s.author)}</span></div>
      </div>
    </div>`).join('');

  document.getElementById('slides').innerHTML = html;
}

function renderSlideBody(s){
  function pill(t){return t?`<div class="s-pill">${esc(t)}</div>`:''}
  function hl(g,w){return `<div class="s-hl"><span class="g">${esc(g)}</span><span class="w"> ${esc(w)}</span></div>`}
  function body(b){
    b = b||['','',''];
    return `<p class="s-body">${esc(b[0])}<mark>${esc(b[1])}</mark>${esc(b[2])}</p>`
  }

  if(s.type==='cover') return `
    ${pill(s.pill)}
    ${hl(s.headline_green, s.headline_white)}
    ${body(s.body)}`;

  if(s.type==='3col') return `
    ${pill(s.pill)}
    ${hl(s.headline_green, s.headline_white)}
    <div class="s-cols">${(s.cols||[]).map(c=>`<div class="s-col">${esc(c)}</div>`).join('')}</div>`;

  if(s.type==='image_2col') return `
    <div class="s-img">[ AI image ]</div>
    ${hl(s.headline_green, s.headline_white)}
    <div class="s-two">${body(s.body_left)}${body(s.body_right)}</div>`;

  if(s.type==='cta') return `
    <div class="s-img">[ AI image ]</div>
    ${pill(s.pill)}
    ${hl(s.headline_green, s.headline_white)}
    ${body(s.body)}`;

  if(s.type==='closing') return `
    <div class="s-closing"><span class="g">${esc(s.closing_green)}</span><span class="w"> ${esc(s.closing_white)}</span></div>
    ${body(s.body)}
    <div class="s-pills-row">${(s.pills||[]).map(p=>pill(p)).join('')}</div>`;

  return '<p style="color:#333;font-size:11px">preview unavailable</p>';
}

async function openFile(path){
  const r = await fetch('/api/open',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({path})});
  const d = await r.json();
  if(!d.ok) alert('Could not open: '+d.error);
}

loadStories();
</script>
</body>
</html>"""

# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--port", type=int, default=5050)
    parser.add_argument("--no-browser", action="store_true")
    args = parser.parse_args()

    if not args.no_browser:
        def _open():
            import time; time.sleep(1.2)
            webbrowser.open(f"http://localhost:{args.port}")
        threading.Thread(target=_open, daemon=True).start()

    print(f"\n  Carousel Preview  →  http://localhost:{args.port}\n  Ctrl+C to stop\n")
    app.run(host="0.0.0.0", port=args.port, debug=False)

if __name__ == "__main__":
    main()
