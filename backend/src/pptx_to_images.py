"""
pptx_to_images.py
-----------------
Convert a carousel PPTX file into PNG images (one per slide).

Method 1: LibreOffice headless (best quality, requires LibreOffice installed)
Method 2: python-pptx + Pillow render (fallback, no external deps)

Usage:
    python pptx_to_images.py path/to/carousel.pptx output/dir/
"""

import os
import sys
import subprocess
import shutil
import glob
from pathlib import Path

try:
    from PIL import Image, ImageDraw, ImageFont
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# Instagram square format: 1080x1080
SLIDE_WIDTH  = 1080
SLIDE_HEIGHT = 1080


def convert_pptx_to_images(pptx_path: str, output_dir: str) -> list:
    """
    Convert a PPTX file into a list of PNG image paths.
    Returns list of absolute paths in slide order.
    """
    pptx_path  = os.path.abspath(pptx_path)
    output_dir = os.path.abspath(output_dir)
    os.makedirs(output_dir, exist_ok=True)

    if not os.path.exists(pptx_path):
        print(f"[pptx_to_images] File not found: {pptx_path}")
        return []

    # Method 1: LibreOffice
    paths = _try_libreoffice(pptx_path, output_dir)
    if paths:
        print(f"[pptx_to_images] Converted {len(paths)} slides via LibreOffice.")
        return paths

    # Method 2: python-pptx + Pillow text render
    if PPTX_AVAILABLE and PIL_AVAILABLE:
        paths = _render_with_pptx_pillow(pptx_path, output_dir)
        if paths:
            print(f"[pptx_to_images] Rendered {len(paths)} slides via Pillow fallback.")
            return paths

    print("[pptx_to_images] ERROR: Could not convert PPTX. Install LibreOffice or python-pptx + Pillow.")
    return []


def _try_libreoffice(pptx_path: str, output_dir: str) -> list:
    """Try converting using LibreOffice headless."""
    soffice = _find_libreoffice()
    if not soffice:
        return []

    try:
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "png", "--outdir", output_dir, pptx_path],
            capture_output=True, text=True, timeout=120
        )
        if result.returncode != 0:
            print(f"[LibreOffice] Error: {result.stderr[:300]}")
            return []

        # Rename and sort output files
        stem = Path(pptx_path).stem
        png_files = sorted(glob.glob(os.path.join(output_dir, f"{stem}*.png")))
        if not png_files:
            # LibreOffice may name them slide1.png, slide2.png etc.
            png_files = sorted(glob.glob(os.path.join(output_dir, "*.png")))

        # Rename to slide_01.png, slide_02.png ...
        renamed = []
        for i, src in enumerate(png_files):
            dst = os.path.join(output_dir, f"slide_{i+1:02d}.png")
            if src != dst:
                os.rename(src, dst)
            renamed.append(dst)

        return renamed

    except (subprocess.TimeoutExpired, FileNotFoundError, Exception) as e:
        print(f"[LibreOffice] Failed: {e}")
        return []


def _find_libreoffice() -> str:
    """Locate the soffice / libreoffice binary."""
    candidates = [
        "soffice",
        "libreoffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/usr/bin/soffice",
        "/usr/bin/libreoffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ]
    for c in candidates:
        if shutil.which(c):
            return c
        if os.path.exists(c):
            return c
    return None


def _render_with_pptx_pillow(pptx_path: str, output_dir: str) -> list:
    """
    Fallback: render each slide as a 1080x1080 PNG using python-pptx text extraction + Pillow.
    This is a simplified text-only render — it won't replicate full design fidelity
    but produces usable images with the slide text content on a dark background.
    """
    prs = Presentation(pptx_path)
    rendered = []

    # Dark brand colours
    BG_COLOR   = (12, 12, 18)
    TEXT_COLOR = (255, 255, 255)
    ACC_COLOR  = (196, 255, 0)   # lime green

    for slide_idx, slide in enumerate(prs.slides):
        img = Image.new("RGB", (SLIDE_WIDTH, SLIDE_HEIGHT), BG_COLOR)
        draw = ImageDraw.Draw(img)

        # Try to load a font; fall back to default
        try:
            title_font = ImageFont.truetype("arial.ttf", 54)
            body_font  = ImageFont.truetype("arial.ttf", 32)
            tag_font   = ImageFont.truetype("arial.ttf", 24)
        except (IOError, OSError):
            title_font = ImageFont.load_default()
            body_font  = ImageFont.load_default()
            tag_font   = ImageFont.load_default()

        # Extract all text from slide
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)

        # Draw green accent bar at top
        draw.rectangle([(0, 0), (SLIDE_WIDTH, 8)], fill=ACC_COLOR)

        # Draw AIPulse branding
        draw.text((40, 30), "⚡ AIPulse", font=tag_font, fill=ACC_COLOR)

        # Draw up to 8 lines of text
        y = 120
        for i, text in enumerate(texts[:8]):
            font  = title_font if i < 2 else body_font
            color = ACC_COLOR  if i == 0 else TEXT_COLOR
            # Word-wrap at ~35 chars
            wrapped = _wrap_text(text, 35)
            for line in wrapped:
                draw.text((40, y), line, font=font, fill=color)
                y += 60 if font == title_font else 44
            y += 10

        # Bottom accent bar
        draw.rectangle([(0, SLIDE_HEIGHT - 8), (SLIDE_WIDTH, SLIDE_HEIGHT)], fill=ACC_COLOR)

        # Slide number indicator dots
        total = len(prs.slides)
        dot_size = 10
        dot_spacing = 20
        start_x = (SLIDE_WIDTH - (total * dot_spacing)) // 2
        for d in range(total):
            color = ACC_COLOR if d == slide_idx else (80, 80, 80)
            draw.ellipse(
                [(start_x + d * dot_spacing, SLIDE_HEIGHT - 35),
                 (start_x + d * dot_spacing + dot_size, SLIDE_HEIGHT - 35 + dot_size)],
                fill=color
            )

        out_path = os.path.join(output_dir, f"slide_{slide_idx+1:02d}.png")
        img.save(out_path, "PNG")
        rendered.append(out_path)

    return rendered


def _wrap_text(text: str, max_chars: int) -> list:
    """Simple word-wrap."""
    words  = text.split()
    lines  = []
    current = ""
    for word in words:
        if len(current) + len(word) + 1 <= max_chars:
            current = (current + " " + word).strip()
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines or [""]


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python pptx_to_images.py <path/to/carousel.pptx> [output_dir]")
        sys.exit(1)

    pptx = sys.argv[1]
    out  = sys.argv[2] if len(sys.argv) > 2 else os.path.join(os.path.dirname(pptx), "slide_images")
    paths = convert_pptx_to_images(pptx, out)
    print(f"Generated {len(paths)} images:")
    for p in paths:
        print(f"  {p}")
